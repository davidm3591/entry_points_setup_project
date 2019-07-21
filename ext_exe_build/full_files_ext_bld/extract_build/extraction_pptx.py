from pptx import Presentation
import glob
from tkinter import *
from tkinter import messagebox
import re
import sys
from extraction_out import *
########################################################################
# Extraction Tool
# Edgenuity, Inc.
########################################################################
# Sub-module
# Code kicked off from extraction_tool.py by calling get_file()
# - 07/04/2019 drm
########################################################################


file_extension_pattern = re.compile(r'pptx')

no_file_title = "FILE MISSING"
no_file_msg = "Please make sure there is a valid file in the input folder."
file_format_error_msg = f'FILE ERROR - Problem File:  '
ext_msg = "File must be a PowerPoint with a file extension type \"pptx\""
ext_msg += "\nCorrect file format: 1111-02-03.pptx"

file_number_msg = "File name must be formatted as 1234-05-06 or 12345-06-07."
file_number_msg += " 12345-06-07.\n\n"
file_number_msg += "    Make sure the file name is formated as:\n\tMBID: 4"
file_number_msg += " or 5 numerical digits\n\tUnit: 2 numerical digits"
file_number_msg += "\n\tLesson: 2 numerical digits"

dash_msg = "File separator type must be a dash \"-\""

time_format_title = "FRAME TIME FORMAT ERROR"
time_format_msg = '''
1. Remove references to seconds and or ":".
     Incorrect: 0:30 or :30 or 30 sec or seconds.
2. Time format must be labeled minutes, minute, or min
     and consist of a decimal minute or minutes.
     Correct: 1.5 or 0.5 minutes, minute, or min.
'''


def get_file():
    """Check for file presence and correct format.
        if the file is there and is the correct
        type and format, load file.
    """
    filename_with_path = glob.glob("C:\\e2020_tool_data\\input\\*.*")

    # This will catch no file in folder condition only
    if not filename_with_path:
        messagebox.showerror(no_file_title, no_file_msg)
        sys.exit(None)

    filename_path = []
    filename_path.append(filename_with_path[0])
    path_string = filename_path[0]

    path_split = path_string.split("\\")
    file = path_split[-1]
    filename = file.split(".")[0]

    separator_chk = re.compile(r'\w*-\w*-\w*')
    dash_sep = separator_chk.findall(filename)

    invalid_file_title = f"ERROR: INVALID/MISSING FILE - {filename}"
    invalid_file_msg = """
    There is either no file or an incorrectly formatted file in
     the \"input\" folder (C:/e2020_tool_data/input).

    FILE FORMAT RULES:
    1.  A 4-digit or 5-digit mbid number
         A 2-digit unit number
         A 2-digit lesson number
    2.  The File must be a .pptx type PowerPoint.
    3.  The mbid, unit, and lesson numbers require
         a (mandatory) separator that is a dash.

    EXAMPLES:
       VALID: 1234-56-78.pptx
       VALID: 12345-67-89.pptx

       INVALID: 1234-56-78.ppt (invalid file ext/type)
       INVALID: 1234-5-6.pptx (invalid file unit/lesson
                       format)
       INVALID: 1234_56_78.pptx (invalid file separator)
       INVALID: USHIU7L1.pptx (invalid name characters
                       - must be numeric)
    """
    if dash_sep:
        fname_list = filename.split('-')
        crs = fname_list[0]
        un = fname_list[1]
        lsn = fname_list[2]

        crs_len_chk = len(crs)
        un_len_chk = len(un)
        lsn_len_chk = len(lsn)

        crs_num = crs.isnumeric()
        un_num = un.isnumeric()
        lsn_num = lsn.isnumeric()

        if (crs_len_chk == 4 or crs_len_chk == 5) and un_len_chk == 2 and lsn_len_chk == 2:
            len_chk = True
        else:
            len_chk = False

        if crs_num and un_num and lsn_num:
            numeric_chk = True
        else:
            numeric_chk = False

        if len_chk and numeric_chk:
            valid_file = True
        else:
            valid_file = False
    else:
        messagebox.showerror(invalid_file_title, invalid_file_msg)
        sys.exit(None)

    if valid_file:
        pass
    else:
        messagebox.showerror(invalid_file_title, invalid_file_msg)
        sys.exit(None)

    mbid_length = len(filename)

    # create separate values for mbid, unit, lesson
    if mbid_length == 10:
        course = filename[0:4]
        unit = filename[5:7]
        lesson = filename[8:10]

    if mbid_length == 11:
        course = filename[0:5]
        unit = filename[6:8]
        lesson = filename[9:11]

    # Create globals and assign course, unit, and lesson values
    #    for use outside of get_file(). (Weird variable names to
    #    reduce the the risk overwrite using globals>)
    if mbid_length == 10 or mbid_length == 11:
        global cnum_x1
        global unum_x1
        global lnum_x1
        cnum_x1 = course
        unum_x1 = unit
        lnum_x1 = lesson
    else:
        pass

    # This will catch the not pptx file extension condition
    chk_ext_format = file_extension_pattern.search(path_string)
    if not chk_ext_format:
        messagebox.showerror(file_format_error_msg + file, ext_msg)
        sys.exit(None)

    file_error_messages = []
    msg_out = ""

    def error_messages(error_msg):
        if file_error_messages:
            messagebox.showerror(file_format_error_msg + file, msg_out)
            sys.exit(None)

        elif error_msg == 'time_format_error':
            messagebox.showerror(time_format_title, time_format_msg)
            sys.exit(None)

    if file_error_messages:
        error_count = 1
        for message in file_error_messages:
            msg_out += f"{error_count}. {message}\n"
            error_count += 1
        error_messages(msg_out)
    # Pass "filename_path" to "read_slide" function
    else:
        read_slide(filename_path)


# Regex for FrameChain, Tagging, Timing, and Vocab
# Tagging (e.g. #Frame)
tags = re.compile(r"^#[a-zA-Z0-9\-]+")
# Timing - rebuilt 1/25/2019
time = re.compile(
    r"timing\s*.\s*(\d*\.?\d*)\s*[a-z]\s*", re.IGNORECASE)
# Vocab (e.g. word – this is the definition)
vocab = re.compile(r"(([\w\-]+ ){1,3})–(.+)")

audio = re.compile(
    r"(\bentry\b|\bhint\b|\bhint1\b|\bhint2\b|\bhint3\b|"
    r"\bhint4\b|\bhint5\b|\bhint6\b|\bhint7\b|\bhint8\b|"
    r"\bclick\b|\bclick1\b|\bclick2\b|\bclick3\b|"
    r"\bclick4\b|\bclick5\b|\bclick6\b|\bclick7\b|\bclick8\b|"
    r"\btrans\w+\b|\btrans\w+1\b|\btrans\w+2\b|\btrans\w+3\b|"
    r"\btrans\w+4\b|\btrans\w+5\b|\btrans\w+6\b|\btrans\w+7\b|"
    r"\btrans\w+8\b|\bexit\b)\sAudio:\s(.+)", re.IGNORECASE)

# Regex to check for ":" separator
unit_pattern_1 = re.compile(r'(Timing\s*.\s*\d*:\d*.)', re.I)
# Regex to check Timing units as 'seconds', 'sec', etc.
unit_pattern_2 = re.compile(
    r'(Timing\s*.\s*\d*.\d*\s*)(sec|second|seconds)', re.I)


def read_slide(filename_path):
    """Read and process PPT note content by slide."""
    prs = Presentation(*filename_path)

    slide_num = 0
    activity_num = 0
    activity_nums = []
    fc_activity = ""
    frame_num = ""

    file_name = ""
    slide_type = ""
    frame_name = ""
    activity_name = ""
    slide_layout = ""
    frame_times = []
    time_key_list = []
    time_key = ""
    frame_directory = []
    vo_output_data = []
    ppt_slide_num = 0
    ext_sheet_data = []
    vid_qa_sheet_data = []
    vo_data_summary = []

    # ########################################
    # Variables and list for basic info
    basic_info = []
    basic_vid_qa = []
    course_chk = False
    unit_chk = False
    Lesson_num_chk = False
    Order_num_chk = False
    lesson_title_chk = False
    # slide_count = 0
    # ########################################

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            # Get basic info for course, unit, lesson, lesson num
            for paragraph in shape.text_frame.paragraphs:
                if lesson_title_chk is False:
                    lesson_title = slide.shapes[2].text
                    lesson_title = f"Lesson Title: {lesson_title}"
                    lesson_title_chk = True

                if course_chk is False:
                    if 'Course:' in paragraph.text:
                        basic_info.append(paragraph.text)
                        basic_info.append(lesson_title)
                        basic_vid_qa.append(lesson_title)
                        course_chk = True

                if unit_chk is False:
                    if 'Unit:' in paragraph.text:
                        basic_info.append(paragraph.text)
                        basic_vid_qa.append(paragraph.text)
                        unit_chk = True

                if Lesson_num_chk is False:
                    if 'Lesson #:' in paragraph.text:
                        basic_info.append(paragraph.text)
                        Lesson_num_chk = True

        ###############################################################
        slide_num += 1
        tags_result = []
        time_result = 0.0
        audio_results = {}
        vocab_results = {}

        # Insert code for base data here?
        # Ref: extract_pptx_text_3_base_info_code.py

        if slide.has_notes_slide:
            # Get notes from PPTX slides
            # https://python-pptx.readthedocs.io/en/latest/user/notes.html
            for line in slide.notes_slide.notes_text_frame.text.split("\n"):

                # Match for Vocab
                v = vocab.match(line.strip())
                if v is not None:
                    vocab_results[v[2]] = v[3]

                # Match for Tags
                t = tags.match(line)
                if t is not None:
                    tags_result.append(t[0])

                # Check units for : or sec(onds)
                unit_matches_1 = unit_pattern_1.findall(line)
                unit_matches_2 = unit_pattern_2.findall(line)
                if unit_matches_1 or unit_matches_2:
                    messagebox.showerror(
                        time_format_title, time_format_msg)
                    sys.exit(None)

                t_chk = time.match(line)
                if t_chk:
                    ti = time.match(line.strip())
                    if ti is not None:
                        time_result = float(ti[1])

                # Match for Audio
                au = audio.match(line.strip())
                if au is not None:
                    audio_results[au[1]] = au[2]

        # Number the activities activity '#activityname'
        activity_type = ""
        if '#activityname' in tags_result:
            activity_num += 1
            fc_activity = str(activity_num * 5)
            if fc_activity == "5":
                fc_activity = f"0{5}"

            # Activity name for ext sht excel
            activity_name = "ActivityName"
            # Activity type for ext sht excel
            if '#fc' in tags_result:
                activity_type = "fc"
            elif '#ea' in tags_result:
                activity_type = "ea"
            elif '#sw' in tags_result:
                activity_type = "sw"
            elif '#uc' in tags_result:
                activity_type = "uc"
            elif '#ea' in tags_result:
                activity_type = "ea"
        else:
            activity_name = ""

        # Number the frames '#activityname' '#frame'
        if "#frame" not in tags_result and "#activityname" in tags_result:
            frame_count = 1
            frame_num = str(f"0{1}")
        if "#frame" in tags_result:
            frame_count = frame_count + 1
            if frame_count < 10:
                frame_num = str(f"0{frame_count}")
            else:
                frame_num = str(frame_count)

        # Number the frame/slide with complete name
        global cnum_x1
        global unum_x1
        global lnum_x1
        slide_base_num = f"{cnum_x1}-{unum_x1}-{lnum_x1}"

        # check if frame has num - create "frame_name"
        if frame_num:
            frame_name = f"{slide_base_num}-{fc_activity}-{frame_num}"
            # build list for directory structure
            frame_directory.append(frame_name)

        # Create file name with file type appended
        # Video file names
        if "#video" in tags_result and "#anchor" in tags_result:
            file_name = f"{frame_name}-anchor"
        elif "#video" in tags_result and "#instruct" in tags_result:
            file_name = f"{frame_name}-instructional"
        elif "#video" in tags_result and "#fs" in tags_result:
            file_name = f"{frame_name}-fs"
        elif "#video" in tags_result and "#title" in tags_result:
            file_name = f"{frame_name}-title"

        # Audio file names
        elif audio_results:
            file_name = []
            for audio_type in audio_results.keys():
                if audio_type == "Entry":
                    file_name_en = f"{frame_name}-{audio_type.lower()}"
                    file_name.append(file_name_en)
                if audio_type == "Hint":
                    file_name_h = f"{frame_name}-{audio_type.lower()}"
                    file_name.append(file_name_h)
                if audio_type == "Hint1":
                    file_name_h1 = f"{frame_name}-{audio_type.lower()}"
                    file_name.append(file_name_h1)
                if audio_type == "Hint2":
                    file_name_h2 = f"{frame_name}-{audio_type.lower()}"
                    file_name.append(file_name_h2)
                if audio_type == "Hint3":
                    file_name_h3 = f"{frame_name}-{audio_type.lower()}"
                    file_name.append(file_name_h3)
                if audio_type == "Hint4":
                    file_name_h4 = f"{frame_name}-{audio_type.lower()}"
                    file_name.append(file_name_h4)
                if audio_type == "Hint5":
                    file_name_h5 = f"{frame_name}-{audio_type.lower()}"
                    file_name.append(file_name_h5)
                if audio_type == "Hint6":
                    file_name_h6 = f"{frame_name}-{audio_type.lower()}"
                    file_name.append(file_name_h6)
                if audio_type == "Hint7":
                    file_name_h7 = f"{frame_name}-{audio_type.lower()}"
                    file_name.append(file_name_h7)
                if audio_type == "Hint8":
                    file_name_h8 = f"{frame_name}-{audio_type.lower()}"
                    file_name.append(file_name_h8)

                if audio_type == "Click":
                    file_name_c = f"{frame_name}-{audio_type.lower()}"
                    file_name.append(file_name_c)
                if audio_type == "Click1":
                    file_name_c1 = f"{frame_name}-{audio_type.lower()}"
                    file_name.append(file_name_c1)
                if audio_type == "Click2":
                    file_name_c2 = f"{frame_name}-{audio_type.lower()}"
                    file_name.append(file_name_c2)
                if audio_type == "Click3":
                    file_name_c3 = f"{frame_name}-{audio_type.lower()}"
                    file_name.append(file_name_c3)
                if audio_type == "Click4":
                    file_name_c4 = f"{frame_name}-{audio_type.lower()}"
                    file_name.append(file_name_c4)
                if audio_type == "Click5":
                    file_name_c5 = f"{frame_name}-{audio_type.lower()}"
                    file_name.append(file_name_c5)
                if audio_type == "Click6":
                    file_name_c6 = f"{frame_name}-{audio_type.lower()}"
                    file_name.append(file_name_c6)
                if audio_type == "Click7":
                    file_name_c7 = f"{frame_name}-{audio_type.lower()}"
                    file_name.append(file_name_c7)
                if audio_type == "Click8":
                    file_name_c8 = f"{frame_name}-{audio_type.lower()}"
                    file_name.append(file_name_c8)

                if audio_type == "Transition":
                    file_name_t = f"{frame_name}-{audio_type.lower()}"
                    file_name.append(file_name_t)
                if audio_type == "Transition1":
                    file_name_t1 = f"{frame_name}-{audio_type.lower()}"
                    file_name.append(file_name_t1)
                if audio_type == "Transition2":
                    file_name_t2 = f"{frame_name}-{audio_type.lower()}"
                    file_name.append(file_name_t2)
                if audio_type == "Transition3":
                    file_name_t3 = f"{frame_name}-{audio_type.lower()}"
                    file_name.append(file_name_t3)
                if audio_type == "Transition4":
                    file_name_t4 = f"{frame_name}-{audio_type.lower()}"
                    file_name.append(file_name_t4)
                if audio_type == "Transition5":
                    file_name_t5 = f"{frame_name}-{audio_type.lower()}"
                    file_name.append(file_name_t5)
                if audio_type == "Transition6":
                    file_name_t6 = f"{frame_name}-{audio_type.lower()}"
                    file_name.append(file_name_t6)
                if audio_type == "Transition7":
                    file_name_t7 = f"{frame_name}-{audio_type.lower()}"
                    file_name.append(file_name_t7)
                if audio_type == "Transition8":
                    file_name_t8 = f"{frame_name}-{audio_type.lower()}"
                    file_name.append(file_name_t8)
                if audio_type == "Exit":
                    file_name_ex = f"{frame_name}-{audio_type.lower()}"
                    file_name.append(file_name_ex)
        else:
            file_name = ""

        # Test for slide type
        # Needs refactoring: switch/case?
        if "#video" in tags_result:
            slide_type = "video"
        elif "#rb" in tags_result:
            slide_type = "radiobutton"
        elif "#cb" in tags_result:
            slide_type = "checkbox"
        elif "#fib" in tags_result:
            slide_type = "fib"
        elif "#dd" in tags_result:
            slide_type = "dd"
        elif "#html" in tags_result:
            slide_type = "basic"
        elif "#2cat" in tags_result:
            slide_type = "Sort by Two Categories"
        elif "#3cat" in tags_result:
            slide_type = "Sort by Three Categories"
        elif "#rank" in tags_result:
            slide_type = "Ranking"
        elif "#match" in tags_result:
            slide_type = "Matching"
        elif "#int" in tags_result:
            slide_type = "Interactive"
        elif "#te" in tags_result:
            slide_type = "Text Evidence"
        else:
            slide_type = ""

        # Test for slide layout
        # Needs refactoring: switch/case?
        if '#fw' in tags_result:
            slide_layout = "full-width"
        elif '#2media-r' in tags_result:
            slide_layout = "2-up-media-right"
        elif '#2media-l' in tags_result:
            slide_layout = "2-up-media-left"
        elif '#2tab' in tags_result:
            slide_layout = "2-up-tabbed"
        elif '#2cb-sv' in tags_result:
            slide_layout = "2-up-checkbox-survey"
        elif '#4html' in tags_result:
            slide_layout = "4-up"
        elif '#3image' in tags_result:
            slide_layout = "three-images"
        elif '#4image' in tags_result:
            slide_layout = ""
        elif '#2dd-r' in tags_result:
            slide_layout = "2up-dd-right"
        elif '#2dd-l' in tags_result:
            slide_layout = "2up-dd-left"
        elif '#2rb-r' in tags_result:
            slide_layout = "2up-radio-right"
        elif '#2rb-l' in tags_result:
            slide_layout = "2up-radio-left"
        elif '#2cb-r' in tags_result:
            slide_layout = "2up-checkbox-right"
        elif '#2cb-l' in tags_result:
            slide_layout = "2up-checkbox-left"

        # Build the output for the "Extraction Sheet" work sheet
        extraction_row = []
        if fc_activity:
            if activity_type:
                extraction_row.append('ActivityType')
                extraction_row.append(activity_type)
                extraction_row.append(fc_activity)
                extraction_row.append(frame_num)
                if not frame_num:
                    extraction_row.append("")
                else:
                    extraction_row.append(frame_name)
                extraction_row.append(slide_num)
                extraction_row.append(slide_type)
                extraction_row.append(slide_layout)

                if file_name:
                    file_name = ''.join(map(str, file_name))
                    extraction_row.append(file_name)
                    extraction_row.append("")
                    extraction_row.append("")
                    extraction_row.append("")
                else:
                    extraction_row.append("")
                    extraction_row.append("")
                    extraction_row.append("")
                    extraction_row.append("")
            elif file_name:
                extraction_row.append("")
                extraction_row.append("")
                extraction_row.append("")
                extraction_row.append(frame_num)
                extraction_row.append(frame_name)
                extraction_row.append(slide_num)
                extraction_row.append(slide_type)
                extraction_row.append(slide_layout)
                if slide_type == 'video':
                    file_name = ''.join(map(str, file_name))
                    extraction_row.append(file_name)
                    extraction_row.append("")
                    extraction_row.append("")
                    extraction_row.append("")
                else:
                    file_name = '\n'.join(map(str, file_name))
                    extraction_row.append(file_name)
                    extraction_row.append("")
                    extraction_row.append("")
                    extraction_row.append("")
            # ##########################################################
            # Added functionality for format video not new frame
            # ##########################################################
            elif not file_name and slide_type == 'video':
                # if slide_type == 'video':
                extraction_row.append("")
                extraction_row.append("")
                extraction_row.append("")
                extraction_row.append("")
                extraction_row.append("")
                extraction_row.append(slide_num)
                extraction_row.append(slide_type)
                extraction_row.append("")
                extraction_row.append("")
                extraction_row.append("")
                extraction_row.append("")
                extraction_row.append("")
            # ##########################################################
            # ##########################################################
            else:
                extraction_row.append("")
                extraction_row.append("")
                extraction_row.append("")
                extraction_row.append("")
                extraction_row.append("")
                extraction_row.append(slide_num)
                extraction_row.append("")
                extraction_row.append("")
                extraction_row.append("")
                extraction_row.append("")
                extraction_row.append("")
                extraction_row.append("")

            ext_sheet_data.append(extraction_row)

        # Build Video QA workheet output
        video_row = []
        if "#video" in tags_result:
            if "anchor" in file_name:
                video_row.append(fc_activity)
                video_row.append(frame_num)
                video_row.append(slide_num)
                video_row.append(file_name)
                video_row.append(slide_type)
                video_row.append("")
                video_row.append("")
                video_row.append("")
                video_row.append("")

            else:
                video_row.append(fc_activity)
                video_row.append(frame_num)
                video_row.append(slide_num)
                video_row.append(file_name)
                video_row.append(slide_type)
                video_row.append("")
                video_row.append("")
                video_row.append("")
                video_row.append("")

            vid_qa_sheet_data.append(video_row)

        # Build VO workheet output
        vo_row = tuple()
        if audio_results:
            for k, v in audio_results.items():
                vo_row = ()
                if k == 'Entry':
                    vo_framename = f'{frame_name}-{k.lower()}'
                    vo_row = (vo_framename, v)
                if k == 'Hint':
                    vo_framename = f'{frame_name}-{k.lower()}'
                    vo_row = (vo_framename, v)
                if k == 'Hint1':
                    vo_framename = f'{frame_name}-{k.lower()}'
                    vo_row = (vo_framename, v)
                if k == 'Hint2':
                    vo_framename = f'{frame_name}-{k.lower()}'
                    vo_row = (vo_framename, v)
                if k == 'Hint3':
                    vo_framename = f'{frame_name}-{k.lower()}'
                    vo_row = (vo_framename, v)
                if k == 'Hint4':
                    vo_framename = f'{frame_name}-{k.lower()}'
                    vo_row = (vo_framename, v)
                if k == 'Hint5':
                    vo_framename = f'{frame_name}-{k.lower()}'
                    vo_row = (vo_framename, v)
                if k == 'Hint6':
                    vo_framename = f'{frame_name}-{k.lower()}'
                    vo_row = (vo_framename, v)
                if k == 'Hint7':
                    vo_framename = f'{frame_name}-{k.lower()}'
                    vo_row = (vo_framename, v)
                if k == 'Hint8':
                    vo_framename = f'{frame_name}-{k.lower()}'
                    vo_row = (vo_framename, v)

                if k == 'Click':
                    vo_framename = f'{frame_name}-{k.lower()}'
                    vo_row = (vo_framename, v)
                if k == 'Click1':
                    vo_framename = f'{frame_name}-{k.lower()}'
                    vo_row = (vo_framename, v)
                if k == 'Click2':
                    vo_framename = f'{frame_name}-{k.lower()}'
                    vo_row = (vo_framename, v)
                if k == 'Click3':
                    vo_framename = f'{frame_name}-{k.lower()}'
                    vo_row = (vo_framename, v)
                if k == 'Click4':
                    vo_framename = f'{frame_name}-{k.lower()}'
                    vo_row = (vo_framename, v)
                if k == 'Click5':
                    vo_framename = f'{frame_name}-{k.lower()}'
                    vo_row = (vo_framename, v)
                if k == 'Click6':
                    vo_framename = f'{frame_name}-{k.lower()}'
                    vo_row = (vo_framename, v)
                if k == 'Click7':
                    vo_framename = f'{frame_name}-{k.lower()}'
                    vo_row = (vo_framename, v)
                if k == 'Click8':
                    vo_framename = f'{frame_name}-{k.lower()}'
                    vo_row = (vo_framename, v)

                if k == 'Transition':
                    vo_framename = f'{frame_name}-{k.lower()}'
                    vo_row = (vo_framename, v)
                if k == 'Transition1':
                    vo_framename = f'{frame_name}-{k.lower()}'
                    vo_row = (vo_framename, v)
                if k == 'Transition2':
                    vo_framename = f'{frame_name}-{k.lower()}'
                    vo_row = (vo_framename, v)
                if k == 'Transition3':
                    vo_framename = f'{frame_name}-{k.lower()}'
                    vo_row = (vo_framename, v)
                if k == 'Transition4':
                    vo_framename = f'{frame_name}-{k.lower()}'
                    vo_row = (vo_framename, v)
                if k == 'Transition5':
                    vo_framename = f'{frame_name}-{k.lower()}'
                    vo_row = (vo_framename, v)
                if k == 'Transition6':
                    vo_framename = f'{frame_name}-{k.lower()}'
                    vo_row = (vo_framename, v)
                if k == 'Transition7':
                    vo_framename = f'{frame_name}-{k.lower()}'
                    vo_row = (vo_framename, v)
                if k == 'Transition8':
                    vo_framename = f'{frame_name}-{k.lower()}'
                    vo_row = (vo_framename, v)
                if k == 'Exit':
                    vo_framename = f'{frame_name}-{k.lower()}'
                    vo_row = (vo_framename, v)

                vo_output_data.append(vo_row)
            vo_data_summary = list(vo_output_data)

        # Set the time keys by acivity + frame \tnumber
        time_key = f"{fc_activity}-{frame_num}"
        # Collect time keys in a list (use slide # as part of key)
        # Add "->" tag to differentiate unit, lesson from activities
        time_key_list.append(f"{slide_num}->{time_key}")
        # Collect frame times in a list
        frame_times.append(f"{time_result}")

        if vocab_results:
            voc_dict = vocab_results

    # Dirctory structure as a set
    frame_set = sorted(set(frame_directory))
    f_name = f'{cnum_x1}-{unum_x1}-{lnum_x1}'

    # Data for PPTX file copy to output folder
    pptx_file = f'{f_name}.pptx'
    nav_dir = 'C:\\e2020_tool_data\\input'
    src_path_file = f'{nav_dir}\\{pptx_file}'
    dest_path = 'C:\\e2020_tool_data\\output'

    # Function call to pass Excel file data
    excel_data_out(vid_qa_sheet_data, ext_sheet_data,
                   vo_data_summary, f_name, voc_dict,
                   frame_set, time_key_list, frame_times,
                   basic_info, basic_vid_qa)

    # Function call to build extraction directory structure
    directory_out(frame_set, f_name)

    # Function call to copy pptx file to output folder
    copy_pptx(nav_dir, src_path_file, dest_path)